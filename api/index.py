import io
import json
import base64
import os
import asyncio
from fastapi import FastAPI, UploadFile, File
from fastapi.responses import HTMLResponse, Response
from pydantic import BaseModel
import PIL.Image
from pptx import Presentation
from pptx.util import Inches, Pt
from anthropic import Anthropic

app = FastAPI()

# Vercel Environment Variables ထဲကနေ Claude API Key ကို လှမ်းယူပါမယ်
API_KEY = os.getenv("ANTHROPIC_API_KEY", "YOUR_CLAUDE_API_KEY_HERE")
client = Anthropic(api_key=API_KEY)

html_content = """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Progress Report Generator</title>
    <script src="https://cdn.tailwindcss.com"></script>
</head>
<body class="bg-slate-50 min-h-screen p-4 md:p-8">
    <div class="max-w-4xl mx-auto bg-white p-6 md:p-10 rounded-xl shadow-lg border border-slate-100">
        <h1 class="text-2xl md:text-3xl font-bold text-slate-800 mb-2 text-center">Site Progress PPTX Generator</h1>
        <p class="text-center text-slate-500 mb-8 text-sm">Powered by Claude 3.5 Vision - Structural, Architectural & M&E works</p>
        
        <div class="border-4 border-dashed border-indigo-200 bg-indigo-50 p-8 md:p-12 text-center rounded-xl cursor-pointer hover:bg-indigo-100 transition" id="drop-zone">
            <input type="file" id="file-input" multiple class="hidden" accept="image/*">
            <svg class="mx-auto h-12 w-12 text-indigo-400 mb-3" fill="none" stroke="currentColor" viewBox="0 0 48 48"><path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M28 8H12a4 4 0 00-4 4v20m32-12v8m0 0v8a4 4 0 01-4 4H12a4 4 0 01-4-4v-4m32-4l-3.172-3.172a4 4 0 00-5.656 0L28 28M8 32l9.172-9.172a4 4 0 015.656 0L28 28m0 0l4 4m4-24h8m-4-4v8m-12 4h.02" /></svg>
            <p class="text-lg text-slate-700 font-medium">Click or Drag & Drop site photos here</p>
            <p class="text-sm text-slate-500 mt-2" id="file-count">0 files selected</p>
        </div>

        <button id="generate-btn" class="w-full mt-6 bg-indigo-600 text-white font-bold py-4 px-4 rounded-lg hover:bg-indigo-700 shadow-md hidden transition">
            Start Claude AI Analysis & Generate PPTX
        </button>

        <div id="progress-container" class="mt-8 hidden">
            <div class="flex justify-between items-center mb-2">
                <p id="status-text" class="text-sm font-semibold text-slate-700">Initializing Claude Engine...</p>
                <p id="percentage-text" class="text-sm font-bold text-indigo-600">0%</p>
            </div>
            <div class="w-full bg-slate-200 rounded-full h-3">
                <div id="progress-bar" class="bg-indigo-500 h-3 rounded-full transition-all duration-500" style="width: 0%"></div>
            </div>
        </div>

        <div id="download-container" class="mt-8 text-center hidden p-6 border-2 border-green-200 bg-green-50 rounded-xl">
            <h2 class="text-xl font-bold text-green-700 mb-2">✅ PowerPoint is Ready!</h2>
            <p class="text-sm text-green-600 mb-4">Photos have been successfully structured and captioned by Claude.</p>
            <button id="download-btn" class="bg-green-600 text-white font-bold py-3 px-8 rounded-lg hover:bg-green-700 shadow-lg transition">Download .PPTX File</button>
        </div>
    </div>

    <script>
        const dropZone = document.getElementById('drop-zone');
        const fileInput = document.getElementById('file-input');
        const fileCount = document.getElementById('file-count');
        const generateBtn = document.getElementById('generate-btn');
        const progressContainer = document.getElementById('progress-container');
        const progressBar = document.getElementById('progress-bar');
        const statusText = document.getElementById('status-text');
        const percentageText = document.getElementById('percentage-text');
        const downloadContainer = document.getElementById('download-container');
        const downloadBtn = document.getElementById('download-btn');

        let selectedFiles = [];
        let pptBlobUrl = null;

        dropZone.addEventListener('click', () => fileInput.click());
        fileInput.addEventListener('change', (e) => {
            selectedFiles = e.target.files;
            fileCount.innerText = `${selectedFiles.length} site photos selected`;
            if (selectedFiles.length > 0) {
                generateBtn.classList.remove('hidden');
                downloadContainer.classList.add('hidden');
            }
        });

        generateBtn.addEventListener('click', async () => {
            generateBtn.classList.add('hidden');
            progressContainer.classList.remove('hidden');
            let aiResults = [];
            
            for (let i = 0; i < selectedFiles.length; i++) {
                statusText.innerText = `Analyzing [${i + 1}/${selectedFiles.length}]: ${selectedFiles[i].name}`;
                let formData = new FormData();
                formData.append("file", selectedFiles[i]);

                try {
                    let response = await fetch('/api/analyze', { method: 'POST', body: formData });
                    let data = await response.json();
                    aiResults.push(data);
                    
                    let percentage = Math.round(((i + 1) / selectedFiles.length) * 100);
                    progressBar.style.width = `${percentage}%`;
                    percentageText.innerText = `${percentage}%`;

                    if (i < selectedFiles.length - 1) {
                        statusText.innerText = `Cooling down API...`;
                        await new Promise(resolve => setTimeout(resolve, 3500));
                    }
                } catch (error) {
                    console.error("Analysis Error:", error);
                }
            }

            statusText.innerText = "Structuring formatting for PowerPoint slides...";
            
            let pptResponse = await fetch('/api/generate_ppt', {
                method: 'POST',
                headers: { 'Content-Type': 'application/json' },
                body: JSON.stringify({ photos: aiResults })
            });

            if (pptResponse.ok) {
                let blob = await pptResponse.blob();
                pptBlobUrl = window.URL.createObjectURL(blob);
                progressContainer.classList.add('hidden');
                downloadContainer.classList.remove('hidden');
            }
        });

        downloadBtn.addEventListener('click', () => {
            if (pptBlobUrl) {
                let a = document.createElement('a');
                a.href = pptBlobUrl;
                a.download = "PYAE_PHYO_AUNG_Progress_Report.pptx";
                document.body.appendChild(a);
                a.click();
                a.remove();
            }
        });
    </script>
</body>
</html>
"""

@app.get("/")
def read_root():
    return HTMLResponse(content=html_content)

@app.post("/api/analyze")
async def analyze_photo(file: UploadFile = File(...)):
    contents = await file.read()
    img = PIL.Image.open(io.BytesIO(contents))
    if img.mode != 'RGB':
        img = img.convert('RGB')
    
    # Resize in memory (800px is optimal for Claude Vision)
    img.thumbnail((800, 800))
    buffered = io.BytesIO()
    img.save(buffered, format="JPEG", quality=85)
    base64_data = base64.b64encode(buffered.getvalue()).decode("utf-8")

    system_prompt = """
    You are an expert Main Contractor Quantity Surveyor and Site Engineer working on Ministry of Education (MOE) school infrastructure and renovation projects in Singapore.
    Your task is to analyze site progress photos and write a very short, professional progress description.
    
    STRICT RULES:
    1. Categorize the photo strictly into EXACTLY ONE of these 8 scopes:
       - Earthworks & Site Clearance
       - Substructure & Foundation
       - Structural Works (Concrete & Formwork)
       - Architectural & Masonry Works
       - Roofing & Waterproofing
       - M&E (Mechanical & Electrical) Works
       - Interior Finishes (Tiling, Painting, Carpentry)
       - External Works & Landscaping
    2. Focus ONLY on the construction activities, materials, and stage of work (e.g., 'Hacking of existing floor tiles', 'Installation of ceiling grids').
    3. DO NOT describe irrelevant background objects.
    4. Use professional construction and engineering terminology.
    5. Keep it brief (1 to 2 sentences maximum).
    
    Return ONLY a valid JSON object in this format: {"category": "Category Name", "caption": "Your caption here"}
    """
    
    try:
        response = client.messages.create(
            model="claude-3-5-sonnet-20241022",
            max_tokens=300,
            system=system_prompt,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": "image/jpeg",
                                "data": base64_data,
                            }
                        },
                        {
                            "type": "text",
                            "text": "Analyze this progress photo and provide the JSON output."
                        }
                    ]
                }
            ]
        )
        
        response_text = response.content[0].text
        # Clean potential markdown formatting from Claude
        clean_json = response_text.replace("```json", "").replace("
```", "").strip()
        data = json.loads(clean_json)
        
        return {"filename": file.filename, "base64_image": base64_data, "category": data["category"], "caption": data["caption"]}
    except Exception as e:
        return {"filename": file.filename, "base64_image": base64_data, "category": "Uncategorized", "caption": "Progress status under review."}

class PhotoData(BaseModel):
    filename: str
    base64_image: str
    category: str
    caption: str

class PPTXRequest(BaseModel):
    photos: list[PhotoData]

@app.post("/api/generate_ppt")
async def generate_ppt(request: PPTXRequest):
    prs = Presentation()
    
    # Title Slide Creation (Professional Branding)
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "PYAE PHYO AUNG"
    subtitle.text = "Site Progress Report\nMinistry of Education (MOE) Infrastructure Project"

    grouped = {}
    for p in request.photos:
        cat = p.category
        if cat not in grouped:
            grouped[cat] = []
        grouped[cat].append(p)
        
    for cat, photos in grouped.items():
        chunks = [photos[i:i + 4] for i in range(0, len(photos), 4)]
        for chunk in chunks:
            slide = prs.slides.add_slide(prs.slide_layouts[5]) 
            title = slide.shapes.title
            title.text = f"Progress Status: {cat}"
            
            num_photos = len(chunk)
            positions = []
            
            if num_photos == 1:
                positions = [(Inches(2.5), Inches(2), Inches(5))]
            elif num_photos == 2:
                positions = [(Inches(0.5), Inches(2), Inches(4)), (Inches(5.5), Inches(2), Inches(4))]
            else:
                positions = [
                    (Inches(0.5), Inches(1.5), Inches(4)), (Inches(5.5), Inches(1.5), Inches(4)),
                    (Inches(0.5), Inches(4.5), Inches(4)), (Inches(5.5), Inches(4.5), Inches(4))
                ]
                
            for idx, photo in enumerate(chunk):
                left, top, width = positions[idx]
                
                # Decode Base64 back to image stream for PPT
                image_stream = io.BytesIO(base64.b64decode(photo.base64_image))
                slide.shapes.add_picture(image_stream, left, top, width=width)
                
                # Caption Placement
                caption_top = top + Inches(3.2) if num_photos <= 2 else top + Inches(2.6)
                txBox = slide.shapes.add_textbox(left, caption_top, width, Inches(0.5))
                txBox.text_frame.text = photo.caption
                for p in txBox.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(11)

    # Save PPT to Memory
    output_stream = io.BytesIO()
    prs.save(output_stream)
    output_stream.seek(0)
    
    return Response(
        content=output_stream.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=PYAE_PHYO_AUNG_Progress_Report.pptx"}
    )